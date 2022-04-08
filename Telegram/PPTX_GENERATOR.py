from pptx import Presentation
import os

def PPTX_GENERATOR(name, UID, today_date):
    prs = Presentation('Certificate_Template1.pptx') 
    for shape in prs.slides[0].shapes:  # перебираем объекты на слайде
        if (shape.has_text_frame):
            if(shape.text_frame.text == 'Name'):  # ищим поле где написано Name
                shape.text_frame.paragraphs[0].runs[0].text = name  # зписываем туда ФИО
            if(shape.text_frame.text == 'DocumentID'):  # ищим поле где написано DocumentID
                shape.text_frame.paragraphs[0].runs[0].text = UID  # # зписываем туда уникальный идентификатор сертификата
    prs.core_properties.author = "Кванториум Новосибирск"  # указываем создателем Кванториум
    prs.core_properties.title = "Сертификат"  # заголовок файла (ВАЖНО ДЛЯ PDF он в заголовке пишется)

    #  для хранения сертификатов создаём две папки, в которых тоже будут папки
    #  внитри GENERATED_PPTX и GENERATED_PDF будут папки-даты
    os.makedirs('GENERATED_PPTX/{}'.format(today_date), exist_ok=True)  # создаём папку
    os.makedirs('GENERATED_PDF/{}'.format(today_date), exist_ok=True)  # создаём папку
    prs.save('GENERATED_PPTX/'+ today_date + '/' + name + '_' + UID + '.pptx')
    return(name + '_' + UID)  # имя файла для перевода его в PDF
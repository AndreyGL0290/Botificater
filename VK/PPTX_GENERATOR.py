from pptx import Presentation
import os

def PPTX_GENERATOR(name, UID, now):

    prs = Presentation('Certificate_Template.pptx')
    for shape in prs.slides[0].shapes:
        if (shape.has_text_frame):
            if(shape.text_frame.text == 'Name'):
                shape.text_frame.paragraphs[0].runs[0].text = name
            if(shape.text_frame.text == 'DocumentID'):
                shape.text_frame.paragraphs[0].runs[0].text = UID
    prs.core_properties.author = "Кванториум Новосибирск"
    prs.core_properties.title = "Сертификат"

    os.makedirs('GENERATED_PPTX/{}'.format(now), exist_ok=True)  # создаём папку
    os.makedirs('GENERATED_PDF/{}'.format(now), exist_ok=True)  # создаём папку
    prs.save('GENERATED_PPTX/'+ now + '/' + name + '_' + UID + '.pptx')
    print("PPTX_GENERATOR - OK")
    
    return(name + '_' + UID)